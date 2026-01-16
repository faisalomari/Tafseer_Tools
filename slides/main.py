# -*- coding: utf-8 -*-
"""
FAST two-step pipeline + preprocessing crop + tqdm progress bars:
1) Download ayah PNG images directly (parallel) with a progress bar.
2) Preprocess each image (crop/stitch) with a progress bar (except ayah 1).
3) Build a PPTX from processed images with a progress bar.

Install:
  pip install requests python-pptx pillow tqdm
"""

from __future__ import annotations

from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import List, Optional, Tuple

import requests
from pptx import Presentation
from PIL import Image
from tqdm import tqdm
# Add these imports near the top (with the other imports)
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re
from bs4 import BeautifulSoup


# =========================
# 🔒 HARDCODED VARIABLES
# =========================

SORA_NUMBER = 20
START_AYAH = 1
MAX_AYAH = 300                 # safety cap
MAX_WORKERS = 16
TIMEOUT_SEC = 20

RAW_DIR = Path("ayah_pngs")
PROC_DIR = Path("ayah_pngs_processed")
# OUTPUT_PPTX = Path("An-Nahl.pptx")

IMG_URL = "https://surahquran.com/img/ayah/{sora}-{ayah}.png"

# =========================
# 🔒 HARDCODED CROP SETTINGS (Y AXIS)
# keep [0:KEEP_TOP_PX], skip [KEEP_TOP_PX:KEEP_TOP_PX+SKIP_MIDDLE_PX], keep rest
# =========================
KEEP_TOP_PX = 210
SKIP_MIDDLE_PX = 120

# Default white:
SLIDE_BG_RGB = (250, 241, 227)

AYA_PAGE_URL = "https://surahquran.com/aya-{ayah}-sora-{sora}.html"
SORA_ENGLISH_HREF = "https://surahquran.com/sorah-english-{sora}.html"

def sanitize_filename(name: str) -> str:
    name = name.strip()
    name = re.sub(r"\s+", " ", name)
    # Windows-safe: keep letters/numbers/space/_/-
    name = re.sub(r"[^A-Za-z0-9 _-]", "", name)
    return (name.replace(" ", "_") or f"Sora_{SORA_NUMBER}")

def get_surah_english_name(session: requests.Session, sora: int) -> str:
    """
    Fetches an HTML ayah page (aya-1-sora-{sora}.html) and extracts the anchor:
      <a href="https://surahquran.com/sorah-english-{sora}.html">NAME</a>
    Returns NAME (e.g., 'An-Nahl'), fallback 'Sora_{sora}'.
    """
    url = AYA_PAGE_URL.format(ayah=1, sora=sora)
    r = session.get(url, timeout=TIMEOUT_SEC)
    r.raise_for_status()

    soup = BeautifulSoup(r.text, "html.parser")
    target_href = SORA_ENGLISH_HREF.format(sora=sora)

    a = soup.find("a", href=target_href)
    if a and a.get_text(strip=True):
        return a.get_text(strip=True)

    return f"Sora_{sora}"


# =========================
# NETWORK DOWNLOAD
# =========================

def make_session() -> requests.Session:
    s = requests.Session()
    s.headers.update({
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120 Safari/537.36"
        )
    })
    return s


def download_one_png(session: requests.Session, sora: int, ayah: int, out_dir: Path) -> Tuple[int, Optional[Path]]:
    url = IMG_URL.format(sora=sora, ayah=ayah)
    out_path = out_dir / f"{sora}-{ayah}.png"

    try:
        r = session.get(url, timeout=TIMEOUT_SEC)
        if r.status_code != 200 or not r.content:
            return ayah, None

        if not r.content.startswith(b"\x89PNG\r\n\x1a\n"):
            return ayah, None

        out_path.write_bytes(r.content)
        return ayah, out_path

    except Exception:
        return ayah, None


def download_all_pngs(sora: int, start_ayah: int, max_ayah: int, out_dir: Path) -> List[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    session = make_session()

    downloaded: List[Tuple[int, Path]] = []

    batch_size = max(8, MAX_WORKERS * 2)
    ayah = start_ayah
    got_any = False

    outer_pbar = tqdm(total=max_ayah - start_ayah + 1, desc="Downloading PNGs", unit="ayah")

    try:
        while ayah <= max_ayah:
            batch = list(range(ayah, min(max_ayah + 1, ayah + batch_size)))

            with ThreadPoolExecutor(max_workers=MAX_WORKERS) as ex:
                futures = [ex.submit(download_one_png, session, sora, a, out_dir) for a in batch]

                for f in as_completed(futures):
                    a, path = f.result()
                    if path is not None:
                        got_any = True
                        downloaded.append((a, path))
                    outer_pbar.update(1)

            # End detection: if we already got some images, and this whole batch has none, stop.
            if got_any:
                batch_found = sum(1 for a in batch if (out_dir / f"{sora}-{a}.png").exists())
                if batch_found == 0:
                    break

            ayah += batch_size

    finally:
        outer_pbar.close()

    downloaded.sort(key=lambda x: x[0])

    # Return contiguous ayahs from start until first missing
    paths: List[Path] = []
    expected = start_ayah
    dmap = {a: p for a, p in downloaded}
    while expected in dmap:
        paths.append(dmap[expected])
        expected += 1

    if not paths:
        raise RuntimeError("No PNGs downloaded. Check sora number or site availability.")

    return paths


# =========================
# IMAGE PREPROCESSING (CROP/STITCH)
# =========================

def remove_horizontal_band_y(in_path: Path, out_path: Path, keep_top: int, skip_mid: int) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)

    with Image.open(in_path) as im:
        im = im.convert("RGBA")
        w, h = im.size

        y0 = max(0, min(keep_top, h))
        y1 = max(0, min(keep_top + skip_mid, h))

        top_part = im.crop((0, 0, w, y0))
        bottom_part = im.crop((0, y1, w, h))

        new_h = top_part.size[1] + bottom_part.size[1]
        stitched = Image.new("RGBA", (w, new_h), (255, 255, 255, 0))

        stitched.paste(top_part, (0, 0))
        stitched.paste(bottom_part, (0, top_part.size[1]))

        stitched.save(out_path, format="PNG")


def preprocess_all_images(raw_paths: List[Path], proc_dir: Path) -> List[Path]:
    proc_dir.mkdir(parents=True, exist_ok=True)

    out_paths: List[Path] = []

    for p in tqdm(raw_paths, desc="Preprocessing (crop)", unit="img"):
        out_p = proc_dir / p.name

        # filename format: "{sora}-{ayah}.png"
        try:
            ayah_num = int(p.stem.split("-")[1])
        except Exception:
            ayah_num = None

        if ayah_num == 1:
            # Keep original for ayah 1
            out_p.write_bytes(p.read_bytes())
        else:
            remove_horizontal_band_y(
                in_path=p,
                out_path=out_p,
                keep_top=KEEP_TOP_PX,
                skip_mid=SKIP_MIDDLE_PX,
            )

        out_paths.append(out_p)

    return out_paths


# =========================
# PPTX BUILD
# =========================

def add_full_image_slide_from_file(prs: Presentation, img_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 1) Set background color by drawing a full-slide rectangle (behind the image)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    r, g, b = SLIDE_BG_RGB
    bg.fill.fore_color.rgb = RGBColor(r, g, b)
    bg.line.fill.background()  # no border

    # Send rectangle to back so image is on top
    bg.element.getparent().remove(bg.element)
    slide.shapes._spTree.insert(2, bg.element)

    # 2) Add the image, fit while preserving aspect ratio
    with Image.open(img_path) as im:
        w_px, h_px = im.size

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    img_ratio = w_px / h_px
    slide_ratio = slide_w / slide_h

    if img_ratio >= slide_ratio:
        width = slide_w
        height = int(slide_w / img_ratio)
        left = 0
        top = int((slide_h - height) / 2)
    else:
        height = slide_h
        width = int(slide_h * img_ratio)
        top = 0
        left = int((slide_w - width) / 2)

    slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)


def build_pptx_from_pngs(png_paths: List[Path], out_pptx: Path) -> None:
    prs = Presentation()

    for p in tqdm(png_paths, desc="Building PPTX", unit="slide"):
        add_full_image_slide_from_file(prs, p)

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))


# =========================
# RUN
# =========================

def main() -> None:
    session = make_session()

    # Auto name PPTX from the HTML element:
    # <a href="https://surahquran.com/sorah-english-16.html">An-Nahl</a>
    surah_name = sanitize_filename(get_surah_english_name(session, SORA_NUMBER))
    output_pptx = Path(f"{surah_name}.pptx")  # e.g. "An-Nahl.pptx"

    raw_pngs = download_all_pngs(
        sora=SORA_NUMBER,
        start_ayah=START_AYAH,
        max_ayah=MAX_AYAH,
        out_dir=RAW_DIR,
    )
    print(f"Downloaded {len(raw_pngs)} PNGs into: {RAW_DIR}")

    processed_pngs = preprocess_all_images(raw_pngs, PROC_DIR)
    print(f"Processed {len(processed_pngs)} PNGs into: {PROC_DIR}")
    print(f"Crop settings: keep top {KEEP_TOP_PX}px, skip next {SKIP_MIDDLE_PX}px, keep rest")
    print("Note: ayah 1 kept original (no crop).")

    build_pptx_from_pngs(processed_pngs, output_pptx)
    print(f"Saved PPTX: {output_pptx}")



if __name__ == "__main__":
    main()
