# -*- coding: utf-8 -*-
"""
Build PPTX for ALL surahs 1..114 from surahquran.com images, then zip them.

Install:
  pip install requests python-pptx pillow tqdm beautifulsoup4
"""

from __future__ import annotations

import re
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import List, Optional, Tuple

import requests
from bs4 import BeautifulSoup
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from tqdm import tqdm


# =========================
# 🔒 HARDCODED VARIABLES
# =========================

SORA_START = 1
SORA_END = 114

START_AYAH = 1
MAX_AYAH = 300                 # safety cap; auto-stop earlier
MAX_WORKERS = 16
TIMEOUT_SEC = 20

# Output folders
SLIDES_DIR = Path("slides")                    # PPTX output folder
RAW_ROOT = Path("ayah_pngs")                   # raw PNGs root
PROC_ROOT = Path("ayah_pngs_processed")        # processed PNGs root
ZIP_OUT = Path("slides.zip")

# URLs
IMG_URL = "https://surahquran.com/img/ayah/{sora}-{ayah}.png"
AYA_PAGE_URL = "https://surahquran.com/aya-{ayah}-sora-{sora}.html"
SORA_ENGLISH_HREF = "https://surahquran.com/sorah-english-{sora}.html"

# Crop settings (Y axis): keep [0:KEEP_TOP_PX], skip next SKIP_MIDDLE_PX, keep rest
KEEP_TOP_PX = 210
SKIP_MIDDLE_PX = 120

# Slide background color (RGB)
SLIDE_BG_RGB = (250, 241, 227)


# =========================
# UTIL
# =========================

def make_session() -> requests.Session:
    s = requests.Session()
    s.headers.update({
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120 Safari/537.36"
        )
    })
    return s


def sanitize_filename(name: str, fallback: str) -> str:
    name = name.strip()
    name = re.sub(r"\s+", " ", name)
    name = re.sub(r"[^A-Za-z0-9 _-]", "", name)  # Windows-safe
    name = name.replace(" ", "_")
    return name or fallback


def get_surah_english_name(session: requests.Session, sora: int) -> str:
    """
    Fetches https://surahquran.com/aya-1-sora-{sora}.html and extracts:
      <a href="https://surahquran.com/sorah-english-{sora}.html">NAME</a>
    """
    url = AYA_PAGE_URL.format(ayah=1, sora=sora)
    r = session.get(url, timeout=TIMEOUT_SEC)
    r.raise_for_status()

    soup = BeautifulSoup(r.text, "html.parser")
    target_href = SORA_ENGLISH_HREF.format(sora=sora)

    a = soup.find("a", href=target_href)
    if a and a.get_text(strip=True):
        return a.get_text(strip=True)

    return f"Sora_{sora}"


# =========================
# DOWNLOAD
# =========================

def download_one_png(session: requests.Session, sora: int, ayah: int, out_dir: Path) -> Tuple[int, Optional[Path]]:
    url = IMG_URL.format(sora=sora, ayah=ayah)
    out_path = out_dir / f"{sora}-{ayah}.png"

    try:
        r = session.get(url, timeout=TIMEOUT_SEC)
        if r.status_code != 200 or not r.content:
            return ayah, None
        if not r.content.startswith(b"\x89PNG\r\n\x1a\n"):
            return ayah, None

        out_path.write_bytes(r.content)
        return ayah, out_path
    except Exception:
        return ayah, None


def download_all_pngs(session: requests.Session, sora: int, start_ayah: int, max_ayah: int, out_dir: Path) -> List[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)

    downloaded: List[Tuple[int, Path]] = []
    batch_size = max(8, MAX_WORKERS * 2)
    ayah = start_ayah
    got_any = False

    pbar = tqdm(total=max_ayah - start_ayah + 1, desc=f"Sora {sora:03d} - Download", unit="ayah", leave=False)

    try:
        while ayah <= max_ayah:
            batch = list(range(ayah, min(max_ayah + 1, ayah + batch_size)))

            with ThreadPoolExecutor(max_workers=MAX_WORKERS) as ex:
                futures = [ex.submit(download_one_png, session, sora, a, out_dir) for a in batch]
                for f in as_completed(futures):
                    a, path = f.result()
                    if path is not None:
                        got_any = True
                        downloaded.append((a, path))
                    pbar.update(1)

            # End detection: if we already got some images and this whole batch has none → stop
            if got_any:
                batch_found = sum(1 for a in batch if (out_dir / f"{sora}-{a}.png").exists())
                if batch_found == 0:
                    break

            ayah += batch_size
    finally:
        pbar.close()

    downloaded.sort(key=lambda x: x[0])

    # Return contiguous ayahs from start until first missing
    paths: List[Path] = []
    expected = start_ayah
    dmap = {a: p for a, p in downloaded}
    while expected in dmap:
        paths.append(dmap[expected])
        expected += 1

    if not paths:
        raise RuntimeError(f"No PNGs downloaded for sora {sora}.")

    return paths


# =========================
# PREPROCESS
# =========================

def remove_horizontal_band_y(in_path: Path, out_path: Path, keep_top: int, skip_mid: int) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)

    with Image.open(in_path) as im:
        im = im.convert("RGBA")
        w, h = im.size

        y0 = max(0, min(keep_top, h))
        y1 = max(0, min(keep_top + skip_mid, h))

        top_part = im.crop((0, 0, w, y0))
        bottom_part = im.crop((0, y1, w, h))

        new_h = top_part.size[1] + bottom_part.size[1]
        stitched = Image.new("RGBA", (w, new_h), (255, 255, 255, 0))

        stitched.paste(top_part, (0, 0))
        stitched.paste(bottom_part, (0, top_part.size[1]))

        stitched.save(out_path, format="PNG")


def preprocess_all_images(raw_paths: List[Path], proc_dir: Path) -> List[Path]:
    proc_dir.mkdir(parents=True, exist_ok=True)
    out_paths: List[Path] = []

    for p in tqdm(raw_paths, desc="Preprocess", unit="img", leave=False):
        out_p = proc_dir / p.name

        # filename: "{sora}-{ayah}.png"
        try:
            ayah_num = int(p.stem.split("-")[1])
        except Exception:
            ayah_num = None

        if ayah_num == 1:
            out_p.write_bytes(p.read_bytes())  # keep original
        else:
            remove_horizontal_band_y(p, out_p, KEEP_TOP_PX, SKIP_MIDDLE_PX)

        out_paths.append(out_p)

    return out_paths


# =========================
# PPTX
# =========================

def add_full_image_slide_from_file(prs: Presentation, img_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # background rectangle
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    r, g, b = SLIDE_BG_RGB
    bg.fill.fore_color.rgb = RGBColor(r, g, b)
    bg.line.fill.background()

    # send to back
    bg.element.getparent().remove(bg.element)
    slide.shapes._spTree.insert(2, bg.element)

    # image fit
    with Image.open(img_path) as im:
        w_px, h_px = im.size

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    img_ratio = w_px / h_px
    slide_ratio = slide_w / slide_h

    if img_ratio >= slide_ratio:
        width = slide_w
        height = int(slide_w / img_ratio)
        left = 0
        top = int((slide_h - height) / 2)
    else:
        height = slide_h
        width = int(slide_h * img_ratio)
        top = 0
        left = int((slide_w - width) / 2)

    slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)


def build_pptx_from_pngs(png_paths: List[Path], out_pptx: Path) -> None:
    prs = Presentation()
    for p in tqdm(png_paths, desc="Build PPTX", unit="slide", leave=False):
        add_full_image_slide_from_file(prs, p)

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))


# =========================
# ZIP
# =========================

def zip_folder(folder: Path, zip_path: Path) -> None:
    if zip_path.exists():
        zip_path.unlink()

    with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for f in sorted(folder.glob("*.pptx")):
            zf.write(f, arcname=f.name)


# =========================
# MAIN (ALL SORA)
# =========================

def main() -> None:
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    RAW_ROOT.mkdir(parents=True, exist_ok=True)
    PROC_ROOT.mkdir(parents=True, exist_ok=True)

    session = make_session()

    all_soras = list(range(SORA_START, SORA_END + 1))
    overall = tqdm(all_soras, desc="All Surahs", unit="sora")

    for sora in overall:
        try:
            # name
            name_raw = get_surah_english_name(session, sora)
            name = sanitize_filename(name_raw, fallback=f"Sora_{sora}")
            out_pptx = SLIDES_DIR / f"{sora:03d}_{name}.pptx"
            overall.set_postfix_str(f"{sora:03d}_{name}")

            # folders per surah
            raw_dir = RAW_ROOT / f"sora_{sora:03d}"
            proc_dir = PROC_ROOT / f"sora_{sora:03d}"

            # download -> preprocess -> pptx
            raw_pngs = download_all_pngs(session, sora, START_AYAH, MAX_AYAH, raw_dir)
            processed_pngs = preprocess_all_images(raw_pngs, proc_dir)
            build_pptx_from_pngs(processed_pngs, out_pptx)

        except Exception as e:
            # continue to next surah
            print(f"\n[WARN] Failed sora {sora}: {e}\n")
            continue

    overall.close()

    # zip the pptx files
    zip_folder(SLIDES_DIR, ZIP_OUT)
    print(f"Done. PPTX files in: {SLIDES_DIR}")
    print(f"ZIP created: {ZIP_OUT}")


if __name__ == "__main__":
    main()
